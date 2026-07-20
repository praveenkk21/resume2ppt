import io
import copy
from pptx import Presentation
from pptx.util import Emu


def extract_shape_inventory(file_bytes: bytes) -> list:
    prs = Presentation(io.BytesIO(file_bytes))
    slides_data = []

    for slide_idx, slide in enumerate(prs.slides):
        shapes_data = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            full_text = "\n".join(
                para.text for para in shape.text_frame.paragraphs
            ).strip()
            shapes_data.append({
                "shape_id": shape.shape_id,
                "shape_name": shape.name,
                "current_text": full_text,
                "bbox": {
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                },
            })
        if shapes_data:
            slides_data.append({"slide_index": slide_idx, "shapes": shapes_data})

    return slides_data


def _replace_text_preserving_format(text_frame, new_text: str):
    new_lines = new_text.split("\n") if new_text else [""]
    paras = text_frame.paragraphs

    for i, line in enumerate(new_lines):
        if i < len(paras):
            para = paras[i]
            runs = para.runs
            if runs:
                # Write the content into the first run, clear the rest
                runs[0].text = line
                for run in runs[1:]:
                    run.text = ""
            else:
                # No runs in this paragraph — add one by manipulating the XML
                from pptx.oxml.ns import qn
                from lxml import etree
                r_elem = etree.SubElement(para._p, qn("a:r"))
                rPr = etree.SubElement(r_elem, qn("a:rPr"), attrib={"lang": "en-US", "dirty": "0"})
                t_elem = etree.SubElement(r_elem, qn("a:t"))
                t_elem.text = line
        else:
            # Need more paragraphs than exist — clone the last one
            last_para = paras[-1]
            new_para = copy.deepcopy(last_para._p)
            # Clear all run texts in the cloned paragraph
            from pptx.oxml.ns import qn
            for r in new_para.findall(qn("a:r")):
                t = r.find(qn("a:t"))
                if t is not None:
                    t.text = ""
            # Set text in first run of cloned para
            runs_in_clone = new_para.findall(qn("a:r"))
            if runs_in_clone:
                t = runs_in_clone[0].find(qn("a:t"))
                if t is not None:
                    t.text = line
            # Append cloned paragraph to txBody
            text_frame._txBody.append(new_para)

    # Clear any remaining original paragraphs that weren't used
    if len(new_lines) < len(paras):
        for para in paras[len(new_lines):]:
            runs = para.runs
            for run in runs:
                run.text = ""
            if not runs:
                # Clear raw text nodes
                from pptx.oxml.ns import qn
                for t in para._p.findall(qn("a:t")):
                    t.text = ""


def apply_mapping(file_bytes: bytes, mapping: dict) -> bytes:
    """
    mapping: {slide_index (str or int): {shape_id (str or int): new_text}}
    """
    prs = Presentation(io.BytesIO(file_bytes))

    for slide_idx_key, shape_map in mapping.items():
        slide_idx = int(slide_idx_key)
        if slide_idx >= len(prs.slides):
            continue
        slide = prs.slides[slide_idx]

        # Build a lookup of shape_id -> new_text
        normalized_map = {str(k): v for k, v in shape_map.items()}

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            sid = str(shape.shape_id)
            if sid not in normalized_map:
                continue
            new_text = normalized_map[sid]
            if new_text is None:
                continue
            _replace_text_preserving_format(shape.text_frame, str(new_text))

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()
